from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 12)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Styling
    for shape in [title, subtitle]:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(248, 250, 252)
                run.font.name = 'Outfit'

def add_comparison_slide(prs, age_title, img_2025, img_2026):
    slide_layout = prs.slide_layouts[5] # Blank with title
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 12)
    
    title = slide.shapes.title
    title.text = f"Comparativo: {age_title}"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(129, 140, 248)
            run.font.name = 'Outfit'
            run.font.size = Pt(32)

    # Images
    # 2025 (Left)
    slide.shapes.add_picture(img_2025, Inches(0.5), Inches(1.5), width=Inches(4.2))
    label_2025 = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(4), Inches(0.5))
    label_2025.text = "Pre-Brote (Abril 2025)"
    
    # 2026 (Right)
    slide.shapes.add_picture(img_2026, Inches(5.3), Inches(1.5), width=Inches(4.2))
    label_2026 = slide.shapes.add_textbox(Inches(5.3), Inches(5.8), Inches(4), Inches(0.5))
    label_2026.text = "Situación Actual (2026)"
    
    for label in [label_2025, label_2026]:
        for paragraph in label.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(148, 163, 184)
                run.font.name = 'Outfit'
                run.font.size = Pt(18)

def create_presentation():
    prs = Presentation()
    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, "Análisis de Cobertura de Vacunación", "Comparativo Sarampión: Durango (2025 - 2026)")
    
    add_comparison_slide(prs, "1 Año (1ra Dosis)", "map_2025_12m.png", "map_2026_12m.png")
    add_comparison_slide(prs, "18 Meses (2da Dosis)", "map_2025_18m.png", "map_2026_18m.png")
    add_comparison_slide(prs, "6 Años (2da Dosis)", "map_2025_6y.png", "map_2026_6y.png")
    
    prs.save("Comparativo_Cobertura_Sarampion.pptx")
    print("Presentation saved as Comparativo_Cobertura_Sarampion.pptx")

if __name__ == "__main__":
    create_presentation()
